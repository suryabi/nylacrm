"""
Knowledge Base — "Ask Nyla" RAG-style assistant.

For an MVP with <50 docs, we skip vector embeddings and stuff all parsed
document text directly into the LLM prompt. This keeps the system portable
(any provider works), simple to maintain, and citation-accurate.

Admin-only document management; all authenticated users can ask questions.
"""
from __future__ import annotations
import io
import os
import uuid
import logging
from datetime import datetime, timezone
from typing import Optional, List, Dict, Any

import httpx
from bs4 import BeautifulSoup
from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Form
from pydantic import BaseModel
from dotenv import load_dotenv

load_dotenv()

from database import db
from deps import get_current_user
from core.tenant import get_current_tenant_id

logger = logging.getLogger(__name__)
router = APIRouter()

# Roles allowed to upload / delete knowledge base documents
ADMIN_ROLES = {"CEO", "System Admin", "Admin"}

# Hard cap on text we send to LLM per question (rough char budget; ~4 chars/token)
MAX_CONTEXT_CHARS = 220_000
MAX_DOC_CHARS = 60_000  # truncate any single doc beyond this

ALLOWED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt", ".md"}


# ════════════════════════════════════════════════════════════════════
# Document parsers — return plain text
# ════════════════════════════════════════════════════════════════════
def _parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join((p.extract_text() or "") for p in reader.pages)


def _parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                out.append(shape.text)
    return "\n".join(out)


def _parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out: List[str] = []
    for ws in wb.worksheets:
        out.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                out.append("\t".join(cells))
    return "\n".join(out)


def _parse_csv_or_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


def _parse_url(url: str) -> str:
    with httpx.Client(timeout=30.0, follow_redirects=True) as client:
        r = client.get(url, headers={"User-Agent": "AskNyla/1.0"})
        r.raise_for_status()
    soup = BeautifulSoup(r.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    title = (soup.title.string.strip() if soup.title and soup.title.string else url)
    text = soup.get_text(separator="\n", strip=True)
    return f"Title: {title}\nURL: {url}\n\n{text}"


def _extract_text(filename: str, data: bytes) -> str:
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == ".pdf":
        return _parse_pdf(data)
    if ext == ".docx":
        return _parse_docx(data)
    if ext == ".pptx":
        return _parse_pptx(data)
    if ext in (".xlsx", ".xls"):
        return _parse_xlsx(data)
    if ext in (".csv", ".txt", ".md"):
        return _parse_csv_or_txt(data)
    raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")


# ════════════════════════════════════════════════════════════════════
# Models
# ════════════════════════════════════════════════════════════════════
class TextDocPayload(BaseModel):
    title: str
    content: str


class UrlDocPayload(BaseModel):
    title: Optional[str] = None
    url: str


class AskPayload(BaseModel):
    question: str
    session_id: Optional[str] = None
    history: Optional[List[Dict[str, str]]] = None  # [{role, content}]


class FeedbackPayload(BaseModel):
    message_id: str
    rating: str  # "up" | "down"
    comment: Optional[str] = None


def _is_admin(user: dict) -> bool:
    return (user or {}).get("role") in ADMIN_ROLES


def _doc_to_dict(doc: dict) -> dict:
    """Strip _id and large content field for list responses."""
    d = {k: v for k, v in doc.items() if k != "_id" and k != "content"}
    d["content_length"] = len(doc.get("content") or "")
    return d


# ════════════════════════════════════════════════════════════════════
# Document management endpoints (admin only for write ops)
# ════════════════════════════════════════════════════════════════════
@router.get("/kb/documents")
async def list_documents(current_user: dict = Depends(get_current_user)):
    tenant_id = get_current_tenant_id()
    docs = await db.kb_documents.find(
        {"tenant_id": tenant_id},
        {"_id": 0, "content": 0},
    ).sort("created_at", -1).to_list(length=200)
    # Add content length
    out = []
    for d in docs:
        d["content_length"] = d.pop("content_length", None) or 0
        out.append(d)
    return {"documents": out}


@router.post("/kb/documents/upload")
async def upload_document(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user),
):
    if not _is_admin(current_user):
        raise HTTPException(status_code=403, detail="Only CEO or System Admin can upload knowledge base documents")
    tenant_id = get_current_tenant_id()

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Empty file")
    if file.size and file.size > 25 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File too large (max 25MB)")

    try:
        text = _extract_text(file.filename or "doc", raw)
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Failed to parse uploaded file")
        raise HTTPException(status_code=400, detail=f"Could not parse file: {e}")

    text = (text or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="No readable text found in file")
    if len(text) > MAX_DOC_CHARS:
        text = text[:MAX_DOC_CHARS]

    doc = {
        "id": str(uuid.uuid4()),
        "tenant_id": tenant_id,
        "title": file.filename,
        "source_type": "file",
        "filename": file.filename,
        "content": text,
        "content_length": len(text),
        "uploaded_by_id": current_user.get("id"),
        "uploaded_by_name": current_user.get("name"),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.kb_documents.insert_one(doc)
    doc.pop("_id", None)
    doc.pop("content", None)
    return doc


@router.post("/kb/documents/text")
async def add_text_document(
    payload: TextDocPayload,
    current_user: dict = Depends(get_current_user),
):
    if not _is_admin(current_user):
        raise HTTPException(status_code=403, detail="Only CEO or System Admin can add knowledge base documents")
    tenant_id = get_current_tenant_id()

    text = (payload.content or "").strip()
    if not text:
        raise HTTPException(status_code=400, detail="Content is empty")
    if len(text) > MAX_DOC_CHARS:
        text = text[:MAX_DOC_CHARS]

    doc = {
        "id": str(uuid.uuid4()),
        "tenant_id": tenant_id,
        "title": payload.title or "Untitled note",
        "source_type": "text",
        "content": text,
        "content_length": len(text),
        "uploaded_by_id": current_user.get("id"),
        "uploaded_by_name": current_user.get("name"),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.kb_documents.insert_one(doc)
    doc.pop("_id", None)
    doc.pop("content", None)
    return doc


@router.post("/kb/documents/url")
async def add_url_document(
    payload: UrlDocPayload,
    current_user: dict = Depends(get_current_user),
):
    if not _is_admin(current_user):
        raise HTTPException(status_code=403, detail="Only CEO or System Admin can add knowledge base documents")
    tenant_id = get_current_tenant_id()

    if not payload.url.startswith(("http://", "https://")):
        raise HTTPException(status_code=400, detail="URL must start with http:// or https://")

    try:
        text = _parse_url(payload.url)
    except Exception as e:
        logger.exception("Failed to fetch URL")
        raise HTTPException(status_code=400, detail=f"Could not fetch URL: {e}")
    if len(text) > MAX_DOC_CHARS:
        text = text[:MAX_DOC_CHARS]

    doc = {
        "id": str(uuid.uuid4()),
        "tenant_id": tenant_id,
        "title": payload.title or payload.url,
        "source_type": "url",
        "url": payload.url,
        "content": text,
        "content_length": len(text),
        "uploaded_by_id": current_user.get("id"),
        "uploaded_by_name": current_user.get("name"),
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.kb_documents.insert_one(doc)
    doc.pop("_id", None)
    doc.pop("content", None)
    return doc


@router.delete("/kb/documents/{doc_id}")
async def delete_document(doc_id: str, current_user: dict = Depends(get_current_user)):
    if not _is_admin(current_user):
        raise HTTPException(status_code=403, detail="Only CEO or System Admin can delete knowledge base documents")
    tenant_id = get_current_tenant_id()
    res = await db.kb_documents.delete_one({"tenant_id": tenant_id, "id": doc_id})
    if not res.deleted_count:
        raise HTTPException(status_code=404, detail="Document not found")
    return {"deleted": True, "id": doc_id}


# ════════════════════════════════════════════════════════════════════
# Ask Nyla — chat endpoint
# ════════════════════════════════════════════════════════════════════
ASK_NYLA_SYSTEM_PROMPT = """You are Nyla, a warm, concise AI sales assistant for the team. You answer questions strictly based on the company's knowledge base documents provided in the user's first message.

RULES:
1. Always cite the source document title in your answer using the format [Doc N] where N is the document number provided.
2. If the answer is NOT in the provided documents, say "I couldn't find this in our knowledge base. You may want to check with a senior team member or update the knowledge base." — DO NOT make up information.
3. Be concise, conversational, and sales-rep friendly.
4. When pricing or specifications are involved, quote them exactly as they appear in the source.
5. If the user asks a follow-up question, use prior context but stay grounded in the documents.

FORMATTING (very important — the UI renders Markdown):
- Use `## ` for section headers and `### ` for sub-section headers (these will appear bold + brand-coloured in the UI).
- Use `- ` for bullet lists; use `**bold**` only for the *key term* at the start of a bullet (e.g. `**Air-Sourced Water**: ...`).
- Keep paragraphs short (2-3 sentences max). Don't wrap entire bullets in bold.
- Inline citations like `[Doc 2]` go at the end of the relevant fact, never inside bold text.
"""


# Small-talk patterns — answered without hitting the RAG pipeline so we don't
# waste tokens or surface random doc snippets when the user just says "hi".
_SMALLTALK_PATTERNS = (
    r"^\s*(hi|hii+|hey+|hello+|yo|hola|namaste|namaskar)\b[\s!.,?]*$",
    r"^\s*(good\s+)?(morning|afternoon|evening|day|night)[\s!.,?]*$",
    r"^\s*how\s+are\s+(you|u|ya)(\s+doing)?[\s?!.,]*$",
    r"^\s*(what's\s+up|whats\s+up|sup|wassup|wsp)[\s?!.,]*$",
    r"^\s*(thanks|thank\s+you|thx|ty|thankyou)[\s?!.,]*$",
    r"^\s*(bye|goodbye|see\s+you|cya)[\s?!.,]*$",
    r"^\s*(who\s+are\s+you|what\s+can\s+you\s+do|help|what\s+do\s+you\s+do)[\s?!.,]*$",
    r"^\s*(ok|okay|cool|nice|great|awesome|got\s+it)[\s?!.,]*$",
)


def _smalltalk_reply(question: str, user_name: str | None = None) -> str | None:
    """Return a canned, friendly reply for greetings/small-talk. None if not smalltalk."""
    import re
    q = (question or "").lower().strip()
    if not q or len(q) > 60:
        return None
    name_part = f", {user_name.split()[0]}" if user_name else ""
    for pat in _SMALLTALK_PATTERNS:
        if re.match(pat, q):
            if "how are you" in q or "how r u" in q:
                return (
                    f"Doing great, thanks for asking{name_part}! 👋\n\n"
                    "I'm ready to help with anything from our knowledge base — pricing, product specs, "
                    "objection handling, sustainability talking points, you name it.\n\n"
                    "What would you like to know?"
                )
            if any(w in q for w in ("thanks", "thank you", "thx", "ty")):
                return "You're welcome! Happy to help — ask me anything else whenever you need."
            if any(w in q for w in ("bye", "goodbye", "see you", "cya")):
                return "Catch you later! 👋 I'll be right here when you need me."
            if "what's up" in q or "whats up" in q or "sup" in q or "wassup" in q or "wsp" in q:
                return f"Hey{name_part}! 👋 Ready to dig into our knowledge base — what would you like to ask?"
            if "who are you" in q or "what can you do" in q or q in ("help",) or "what do you do" in q:
                return (
                    "## I'm Nyla 👋\n\n"
                    "Your AI sales assistant, grounded in this company's knowledge base. I can help with:\n\n"
                    "- **Product details** — specs, SKUs, pricing\n"
                    "- **Sales playbooks** — objection handling, talking points\n"
                    "- **Sustainability story** — air-water sourcing, eco operations\n"
                    "- **Process & policy** — warranties, returns, escalation paths\n\n"
                    "Every answer comes with [Doc N] citations so you can verify the source. What would you like to know?"
                )
            if any(w in q for w in ("ok", "okay", "cool", "nice", "great", "awesome", "got it")):
                return "👍 Let me know if there's anything else you'd like to explore."
            # Default greeting (hi / hello / good morning / etc.)
            return (
                f"Hi{name_part}! 👋 I'm Nyla, your AI sales assistant.\n\n"
                "Ask me anything about our products, pricing, processes, or sustainability story — "
                "I'll answer using our knowledge base with citations."
            )
    return None


def _build_context(docs: List[dict]) -> tuple[str, List[dict]]:
    """Concatenate doc texts with [Doc N] markers; respect MAX_CONTEXT_CHARS budget."""
    pieces: List[str] = []
    citations: List[dict] = []
    used = 0
    for i, d in enumerate(docs, 1):
        header = f"\n\n=== [Doc {i}] {d.get('title', 'Untitled')} ===\n"
        body = d.get("content", "") or ""
        block = header + body
        if used + len(block) > MAX_CONTEXT_CHARS:
            remaining = MAX_CONTEXT_CHARS - used - len(header) - 100
            if remaining <= 0:
                break
            block = header + body[:remaining] + "\n[...truncated...]"
        pieces.append(block)
        citations.append({"index": i, "id": d.get("id"), "title": d.get("title"), "source_type": d.get("source_type")})
        used += len(block)
        if used >= MAX_CONTEXT_CHARS:
            break
    return "".join(pieces), citations


@router.post("/kb/ask")
async def ask_nyla(payload: AskPayload, current_user: dict = Depends(get_current_user)):
    """Answer a question using the tenant's knowledge base documents."""
    tenant_id = get_current_tenant_id()
    question = (payload.question or "").strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question is empty")

    session_id = payload.session_id or str(uuid.uuid4())

    # Fast-path: greetings / small-talk — answer without RAG so we don't
    # surface random doc snippets to "hi" or "how are you".
    canned = _smalltalk_reply(question, current_user.get("name"))
    if canned is not None:
        msg_id = str(uuid.uuid4())
        await db.kb_messages.insert_one({
            "id": msg_id,
            "tenant_id": tenant_id,
            "user_id": current_user.get("id"),
            "user_name": current_user.get("name"),
            "session_id": session_id,
            "question": question,
            "response": canned,
            "citation_count": 0,
            "smalltalk": True,
            "created_at": datetime.now(timezone.utc).isoformat(),
        })
        return {
            "id": msg_id,
            "session_id": session_id,
            "answer": canned,
            "citations": [],
        }

    # Pull all documents for this tenant
    docs = await db.kb_documents.find(
        {"tenant_id": tenant_id},
        {"_id": 0},
    ).sort("created_at", 1).to_list(length=100)

    if not docs:
        raise HTTPException(
            status_code=400,
            detail="The knowledge base is empty. Ask your admin to upload documents first.",
        )

    context_text, citations = _build_context(docs)

    # Build prompt — fold history into a single Gemini call (multi-turn is achieved by
    # prepending recent user messages to the prompt body).
    history_lines = []
    if payload.history:
        for h in payload.history[-6:]:
            role = h.get("role")
            content = (h.get("content") or "").strip()
            if not content:
                continue
            if role == "user":
                history_lines.append(f"Earlier user message: {content}")
            elif role == "assistant":
                history_lines.append(f"Your earlier answer: {content}")

    user_prompt_parts = []
    if history_lines:
        user_prompt_parts.append("Conversation so far:")
        user_prompt_parts.extend(history_lines)
        user_prompt_parts.append("")
    user_prompt_parts.extend([
        "Here is our knowledge base:",
        context_text,
        "",
        f"Question: {question}",
        "",
        "Answer concisely using only the knowledge base above. Cite source documents inline as [Doc N].",
    ])
    user_prompt = "\n".join(user_prompt_parts)

    # Call Gemini directly with the user's own API key.
    try:
        from utils.gemini_helpers import gemini_text
        response = await gemini_text(
            prompt=user_prompt,
            system=ASK_NYLA_SYSTEM_PROMPT,
        )
    except HTTPException:
        raise
    except RuntimeError as e:
        raise HTTPException(status_code=500, detail=str(e))
    except Exception as e:
        logger.exception("LLM call failed")
        raise HTTPException(status_code=502, detail=f"AI service error: {e}")

    # Persist message
    msg_id = str(uuid.uuid4())
    await db.kb_messages.insert_one({
        "id": msg_id,
        "tenant_id": tenant_id,
        "user_id": current_user.get("id"),
        "user_name": current_user.get("name"),
        "session_id": session_id,
        "question": question,
        "response": response,
        "citation_count": len(citations),
        "created_at": datetime.now(timezone.utc).isoformat(),
    })

    return {
        "id": msg_id,
        "session_id": session_id,
        "answer": response,
        "citations": citations,
    }


@router.post("/kb/feedback")
async def submit_feedback(payload: FeedbackPayload, current_user: dict = Depends(get_current_user)):
    if payload.rating not in ("up", "down"):
        raise HTTPException(status_code=400, detail="rating must be 'up' or 'down'")
    tenant_id = get_current_tenant_id()
    await db.kb_feedback.insert_one({
        "id": str(uuid.uuid4()),
        "tenant_id": tenant_id,
        "message_id": payload.message_id,
        "user_id": current_user.get("id"),
        "rating": payload.rating,
        "comment": payload.comment,
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    return {"ok": True}


@router.get("/kb/recent")
async def recent_questions(current_user: dict = Depends(get_current_user)):
    """Recent questions asked by current user (for resume / suggestions)."""
    tenant_id = get_current_tenant_id()
    user_id = current_user.get("id")
    msgs = await db.kb_messages.find(
        {"tenant_id": tenant_id, "user_id": user_id},
        {"_id": 0, "question": 1, "created_at": 1, "session_id": 1, "id": 1},
    ).sort("created_at", -1).to_list(length=10)
    return {"messages": msgs}
